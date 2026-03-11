
from pptx import Presentation

try:
    prs = Presentation(r'd:\Projects\Forage TATA Intern\TASK_DETAILS\TASK_4\Presentation_Template.pptx')
    print(f"Total Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}: Layout Name: {slide.slide_layout.name}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"  - Shape Text: {shape.text[:50]}...")
except Exception as e:
    print(f"Error inspecting PPT: {e}")
