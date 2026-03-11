
from pptx import Presentation
from pptx.util import Inches, Pt

def add_slide(prs, title_text, content_text):
    # Try using layout 1 (usually Title and Content)
    # If not available or suitable, use layout 6 (Blank) and add text boxes
    
    layout = prs.slide_layouts[1] # Typically Title and Content
    try:
        slide = prs.slides.add_slide(layout)
    except:
        layout = prs.slide_layouts[0] # Fallback to Title
        slide = prs.slides.add_slide(layout)

    # Set Title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    else:
        # Create title box manually if layout has no title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True

    # Set Content
    # Try to find a body placeholder
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
    
    if body:
        body.text = content_text
    else:
        # Create content box manually
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = content_text
        p.font.size = Pt(18)

def create_presentation():
    # Load template
    try:
        prs = Presentation(r'd:\Projects\Forage TATA Intern\TASK_DETAILS\TASK_4\Presentation_Template.pptx')
    except:
        prs = Presentation() # Fallback to blank if template fails

    # Slide 1: Title
    # Use layout 0 for title slide usually
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = "Predictive Analytics & Autonomous Collections Strategy"
    
    # Subtitle placeholder is usually idx 1 on title slide
    try:
        slide.placeholders[1].text = "Final Report for Geldium\nTata Data Analysis Virtual Internship"
    except:
        pass # Ignore if no subtitle placeholder

    # Slide 2: Problem Statement
    add_slide(prs, "Business Problem & Objective", 
              "Problem: High delinquency rates and inefficient manual collections processes.\n\n"
              "Objective:\n"
              "1. Develop a predictive model to identify high-risk customers.\n"
              "2. Implement an autonomous Agentic AI system to personalize interventions.\n"
              "3. Ensure regulatory compliance (ECOA, GDPR) and fairness.")

    # Slide 3: Model Methodology
    add_slide(prs, "Model Methodology: Random Forest",
              "Methodology:\n"
              "- Algorithm: Random Forest Classifier (Ensemble Method).\n"
              "- Handling Imbalance: SMOTE (Synthetic Minority Over-sampling Technique).\n"
              "- Features: Credit Utilization, Income, Previous Missed Payments.\n\n"
              "Performance:\n"
              "- Recall: Maximized to detect potential defaulters.\n"
              "- ROC-AUC: Strong discrimination between risk classes.")

    # Slide 4: Operational Strategy
    add_slide(prs, "Operational Strategy: Autonomous Agentic AI",
              "The system moves from static rules to dynamic decisions:\n\n"
              "1. Data Pipeline: Real-time ingestion of transaction & bureau data.\n"
              "2. Decision Engine: Contextual reasoning (e.g., 'Distress' vs 'Habitual').\n"
              "3. Action Layer:\n"
              "   - Low Risk: Monitor only.\n"
              "   - Medium Risk: Automated 'Nudges' (SMS/Email).\n"
              "   - High Risk: Human Intervention (Call/Hardship Offer).")

    # Slide 5: Ethical Governance
    add_slide(prs, "Ethical Governance & Compliance Framework",
              "Ensuring Responsible AI:\n\n"
              "- Explainability: SHAP values provide 'Reason Codes'.\n"
              "- Fairness: Automated Disparate Impact Analysis across Age/Location.\n"
              "- Compliance:\n"
              "  - ECOA: Anti-discrimination monitoring.\n"
              "  - GDPR: Right to explanation.\n"
              "  - Human-in-the-Loop: Review for adverse actions.")

    # Slide 6: Conclusion
    add_slide(prs, "Strategic Impact & ROI",
              "Expected Business Outcomes:\n\n"
              "1. Reduced Delinquency: 10-15% reduction in roll rates.\n"
              "2. Operational Efficiency: 40% reduction in manual call volume.\n"
              "3. Customer Trust: Transparent, supportive engagement.\n\n"
              "Next Steps: Pilot phase with 'Credit Health' SMS campaign.")

    output_path = 'Final_Project_Presentation.pptx'
    prs.save(output_path)
    print(f"Presentation saved as {output_path}")

if __name__ == "__main__":
    create_presentation()
