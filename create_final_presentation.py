
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide_with_title_content(prs, title_text, content_text):
    # Use Layout 1: Title and Content
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    slide.shapes.title.text = title_text
    
    # Content
    # Placeholder 1 is usually the content box
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = content_text
            # Optional: formatting
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
            return slide
            
    # Fallback if no placeholder
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    txBox.text_frame.text = content_text
    return slide

def add_slide_with_table(prs, title_text, data):
    # Use Layout 5: Title Only (usually) to have space for table
    slide_layout = prs.slide_layouts[5] 
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    
    # Add Table
    rows = len(data)
    cols = len(data[0])
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(0.8) # arbitrary

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(4)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            # Format
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(16)
                if r == 0: # Header
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(18)

def create_final_deck():
    try:
        prs = Presentation(r'd:\Projects\Forage TATA Intern\TASK_DETAILS\TASK_4\Presentation_Template.pptx')
    except:
        prs = Presentation() # Fallback

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI-Powered Autonomous Collections Strategy"
    try:
        slide.placeholders[1].text = "Proposed System Design for Geldium\nTata Data Analysis Virtual Internship"
    except:
        pass

    # 2. How the System Works (Workflow)
    content_text = (
        "1. INPUTS (Data Ingestion)\n"
        "   - Real-time Transaction Data\n"
        "   - Credit Bureau Updates\n"
        "   - Customer Interaction History (CRM)\n\n"
        "2. DECISION LOGIC (The Brain)\n"
        "   - Predictive Model: Random Forest Risk Score (0-100%)\n"
        "   - Contextual Rules: Detect 'Distress' (e.g., job loss) vs 'Habitual'\n\n"
        "3. ACTIONS (The Execution)\n"
        "   - Autonomous Nudges: SMS/Email for Medium Risk\n"
        "   - Human Escalation: Call tasks for High Risk\n\n"
        "4. LEARNING LOOP (Feedback)\n"
        "   - Reinforcement Learning: System 'learns' which channel gets a response."
    )
    add_slide_with_title_content(prs, "System Workflow: From Data to Action", content_text)

    # 3. Role of Agentic AI (Table)
    table_data = [
        ["Autonomous AI Tasks", "Human Oversight Required"],
        ["Monitoring Payment Behavior", "Approving Adverse Actions (e.g., Reporting)"],
        ["Sending Low/Medium Risk 'Nudges'", "Handling Complex Disputes & Legal Cases"],
        ["Offering Standard Payment Plans", "Reviewing Bias & Fairness Audits"],
        ["Optimizing Contact Timing (e.g., 6 PM)", "Managing Vulnerable Customers (e.g., Deceased)"]
    ]
    add_slide_with_table(prs, "Role of Agentic AI: Autonomy vs. Oversight", table_data)

    # 4. Responsible AI Guardrails (List)
    guardrails_text = (
        "• FAIRNESS: Nightly 'Disparate Impact' Analysis\n"
        "  - Automated check to ensure no demographic group is targeted >20% more than others.\n\n"
        "• EXPLAINABILITY: SHAP-based 'Reason Codes'\n"
        "  - Every decision comes with a generated explanation (e.g., 'Risk flagged due to Utilization > 80%').\n\n"
        "• COMPLIANCE: Hard-Coded Regulatory Rules\n"
        "  - 'Do Not Contact' time windows enforces legal calling hours.\n"
        "  - 'Circuit Breakers' halt the model if error rates spike.\n\n"
        "• TRANSPARENCY: Customer Appeal Pathway\n"
        "  - Simple mechanism for customers to dispute AI decisions."
    )
    add_slide_with_title_content(prs, "Responsible AI Guardrails", guardrails_text)

    # 5. Expected Business Impact (Impact)
    impact_text = (
        "QUANTITATIVE IMPACT (Business KPIs):\n"
        "• 10-15% Reduction in Delinquency Roll-Rates (Early Intervention)\n"
        "• 40% Reduction in Manual Operational Costs (Automating Low Risk)\n"
        "• 20% Increase in Recovery on Early-Stage Debt\n\n"
        "QUALITATIVE IMPACT (Customer Outcomes):\n"
        "• Improved Customer Experience (Less intrusive, more supportive)\n"
        "• Greater Fairness & Consistency (Removing human bias)\n"
        "• Scalability (Handle 10x volume without hiring 10x staff)"
    )
    add_slide_with_title_content(prs, "Expected Business Impact", impact_text)

    save_path = 'Final_Project_Presentation.pptx'
    prs.save(save_path)
    print(f"Deck saved to {save_path}")

if __name__ == "__main__":
    create_final_deck()
