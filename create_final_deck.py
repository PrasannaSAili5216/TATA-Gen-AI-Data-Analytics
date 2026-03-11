
from pptx import Presentation
from pptx.util import Inches, Pt

def add_content_slide(prs, title, content):
    layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Content placeholder
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
            
    if body:
        body.text = content
        for p in body.text_frame.paragraphs:
            p.font.size = Pt(18)
    else:
        # Fallback
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        txBox.text_frame.text = content
        for p in txBox.text_frame.paragraphs:
            p.font.size = Pt(18)

def add_table_slide(prs, title, data):
    try:
        layout = prs.slide_layouts[5] # Title Only
    except:
        layout = prs.slide_layouts[6] # Blank
    
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    else:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        txBox.text_frame.text = title
        if txBox.text_frame.paragraphs:
            txBox.text_frame.paragraphs[0].font.size = Pt(32)
            txBox.text_frame.paragraphs[0].font.bold = True
            
    rows = len(data)
    cols = len(data[0])
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(9.0)
    height = Inches(0.8)
    
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table
    
    # Set columns width
    table.columns[0].width = Inches(4.5)
    table.columns[1].width = Inches(4.5)
    
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                if r == 0:
                    p.font.bold = True
                    p.font.size = Pt(16)

def create_deck():
    path = r'd:\Projects\Forage TATA Intern\TASK_DETAILS\TASK_4\Presentation_Template.pptx'
    try:
        prs = Presentation(path)
    except:
        prs = Presentation()
        
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = "AI-Powered Autonomous Collections Strategy"
    try:
        slide.placeholders[1].text = "Proposed System Design for Geldium\nTata Data Analysis Virtual Internship"
    except:
        pass
        
    # Slide 2: Workflow
    workflow_text = (
        "1. INPUTS (Real-Time Ingestion)\n"
        "   - Transaction Data & Bureau Updates feeds.\n"
        "   - Customer Interaction History from CRM.\n\n"
        "2. DECISION LOGIC (Agentic Reasoning)\n"
        "   - Predictive Model: Random Forest Risk Score.\n"
        "   - Context Engine: Distinguishes 'Hardship' from 'Avoidance'.\n\n"
        "3. ACTIONS (Autonomous Execution)\n"
        "   - Low/Medium Risk: Automated SMS/Email Nudges.\n"
        "   - High Risk: Escalation to Human Collections Team.\n\n"
        "4. LEARNING LOOP (Continuous Improvement)\n"
        "   - Reinforcement Learning: Optimizes channel & timing based on response."
    )
    add_content_slide(prs, "System Workflow: From Data to Action", workflow_text)
    
    # Slide 3: Agentic AI Table
    t_data = [
        ["Autonomous AI Tasks", "Human Oversight Required"],
        ["Monitoring Payment Behavior", "Approving Adverse Actions (e.g., Reporting)"],
        ["Sending Low/Medium Risk 'Nudges'", "Handling Complex Disputes & Legal Cases"],
        ["Offering Standard Payment Deferrals", "Reviewing Bias & Fairness Audits"],
        ["Optimizing Contact Timing (e.g., 6 PM)", "Managing Vulnerable Customers (Deceased)"]
    ]
    add_table_slide(prs, "Role of Agentic AI: Autonomy vs. Oversight", t_data)
    
    # Slide 4: Guardrails
    guard = (
        "• FAIRNESS: Nightly 'Disparate Impact' Analysis\n"
        "  - Automated check to ensure no demographic group is targeted >20% more than others.\n\n"
        "• EXPLAINABILITY: SHAP-based 'Reason Codes'\n"
        "  - Every decision has a generated explanation (e.g., 'Risk flagged due to Utilization > 80%').\n\n"
        "• COMPLIANCE: Regulatory 'Circuit Breakers'\n"
        "  - Hard-coded rules (e.g., 8 AM - 9 PM contact window) to prevent violations.\n\n"
        "• TRANSPARENCY: Customer Appeal Pathway\n"
        "  - Simple mechanism for customers to dispute AI decisions."
    )
    add_content_slide(prs, "Responsible AI Guardrails", guard)
    
    # Slide 5: Impact
    impact = (
        "QUANTITATIVE IMPACT (Business KPIs):\n"
        "• 10-15% Reduction in Delinquency Roll-Rates.\n"
        "• 40% Reduction in Manual Operational Costs.\n"
        "• 20% Increase in Early-Stage Debt Recovery.\n\n"
        "QUALITATIVE IMPACT (Customer Outcomes):\n"
        "• Improved Customer Experience (Supportive, not intrusive).\n"
        "• Greater Fairness & Consistency (Removing human bias).\n"
        "• Scalability (Handle volume spikes without hiring)."
    )
    add_content_slide(prs, "Expected Business Impact", impact)
    
    out = 'Final_Submission_Deck.pptx'
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == '__main__':
    create_deck()
